#!/usr/bin/env python3
"""Emit a deterministic template/master/layout/placeholder inventory as JSON."""

from __future__ import annotations

import argparse
import hashlib
import json
import zipfile
from pathlib import Path

SCHEMA_VERSION = 1


def enum_name(value: object) -> str:
    return getattr(value, "name", str(value))


def placeholder_record(shape: object) -> dict[str, object]:
    placeholder_format = shape.placeholder_format
    return {
        "idx": placeholder_format.idx,
        "type": enum_name(placeholder_format.type),
        "name": shape.name,
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }


def theme_inventory(path: Path) -> list[dict[str, object]]:
    with zipfile.ZipFile(path) as archive:
        result = []
        for name in sorted(info.filename for info in archive.infolist() if not info.is_dir()):
            if name.startswith("ppt/theme/") and name.endswith(".xml"):
                data = archive.read(name)
                result.append(
                    {"part": name, "size": len(data), "sha256": hashlib.sha256(data).hexdigest()}
                )
        return result


def inspect_template(path: Path) -> dict[str, object]:
    from pptx import Presentation

    presentation = Presentation(path)
    masters = []
    layouts = []
    global_index = 0
    for master_index, master in enumerate(presentation.slide_masters):
        master_layouts = []
        for master_layout_index, layout in enumerate(master.slide_layouts):
            placeholders = sorted(
                (placeholder_record(shape) for shape in layout.placeholders),
                key=lambda item: (int(item["idx"]), str(item["name"])),
            )
            record = {
                "name": layout.name,
                "index": global_index,
                "master_index": master_index,
                "master_layout_index": master_layout_index,
                "placeholders": placeholders,
            }
            layouts.append(record)
            master_layouts.append(global_index)
            global_index += 1
        masters.append(
            {
                "index": master_index,
                "name": getattr(master, "name", "") or "",
                "layout_indices": master_layouts,
                "placeholders": sorted(
                    (placeholder_record(shape) for shape in master.placeholders),
                    key=lambda item: (int(item["idx"]), str(item["name"])),
                ),
            }
        )
    name_counts: dict[str, int] = {}
    for layout in layouts:
        name = str(layout["name"])
        name_counts[name] = name_counts.get(name, 0) + 1
    return {
        "schema_version": SCHEMA_VERSION,
        "ok": True,
        "path": str(path),
        "slide_width": presentation.slide_width,
        "slide_height": presentation.slide_height,
        "masters": masters,
        "master_count": len(masters),
        "layouts": layouts,
        "duplicate_layout_names": sorted(name for name, count in name_counts.items() if count > 1),
        "themes": theme_inventory(path),
    }


def parser() -> argparse.ArgumentParser:
    result = argparse.ArgumentParser(description=__doc__)
    result.add_argument("template", type=Path)
    return result


def main(argv: list[str] | None = None) -> int:
    args = parser().parse_args(argv)
    try:
        payload = inspect_template(args.template.resolve())
    except Exception as exc:  # deterministic boundary for library/package parse failures
        payload = {
            "schema_version": SCHEMA_VERSION,
            "ok": False,
            "error": f"{type(exc).__name__}: {exc}",
        }
        print(json.dumps(payload, indent=2, sort_keys=True))
        return 2
    print(json.dumps(payload, indent=2, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
