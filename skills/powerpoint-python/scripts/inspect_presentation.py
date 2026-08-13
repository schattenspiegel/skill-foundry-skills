#!/usr/bin/env python3
"""Emit a bounded read-only semantic and OOXML inventory for a presentation."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import xml.etree.ElementTree as ET
import zipfile
from collections.abc import Iterable
from pathlib import Path
from typing import Any

SCHEMA_VERSION = 1
DEFAULT_DETAILS_LIMIT = 200
MACRO_CONTENT_TYPE = "application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml"
RISK_PATTERNS: tuple[tuple[str, re.Pattern[str]], ...] = (
    ("vba", re.compile(r"(^|/)vbaProject(?:Signature)?\.bin$", re.I)),
    ("embeddings", re.compile(r"^ppt/embeddings/", re.I)),
    ("media", re.compile(r"^ppt/media/", re.I)),
    ("charts", re.compile(r"^ppt/charts/", re.I)),
    ("diagrams", re.compile(r"^ppt/diagrams/", re.I)),
    ("comments", re.compile(r"^ppt/(comments|people)/", re.I)),
    ("tags", re.compile(r"^ppt/tags/", re.I)),
    ("custom_xml", re.compile(r"^customXml/", re.I)),
    ("signatures", re.compile(r"^_xmlsignatures/", re.I)),
)
KNOWN_PREFIXES = (
    "[Content_Types].xml",
    "_rels/",
    "docProps/",
    "customXml/",
    "ppt/presentation.xml",
    "ppt/presProps.xml",
    "ppt/viewProps.xml",
    "ppt/tableStyles.xml",
    "ppt/commentAuthors.xml",
    "ppt/slideMasters/",
    "ppt/slideLayouts/",
    "ppt/slides/",
    "ppt/theme/",
    "ppt/notesSlides/",
    "ppt/notesMasters/",
    "ppt/handoutMasters/",
    "ppt/charts/",
    "ppt/embeddings/",
    "ppt/media/",
    "ppt/comments/",
    "ppt/people/",
    "ppt/tags/",
    "ppt/diagrams/",
    "ppt/fonts/",
    "ppt/printerSettings/",
    "ppt/vbaProject",
)


def enum_name(value: object) -> str:
    return getattr(value, "name", str(value))


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def stable_digest(items: Iterable[object]) -> str:
    payload = json.dumps(list(items), ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def bounded(items: list[Any], limit: int) -> dict[str, object]:
    return {
        "count": len(items),
        "items": items[:limit],
        "truncated": len(items) > limit,
        "digest": stable_digest(items),
    }


def relationships(
    archive: zipfile.ZipFile, names: list[str]
) -> tuple[list[dict[str, str]], list[dict[str, str]]]:
    records: list[dict[str, str]] = []
    malformed: list[dict[str, str]] = []
    for name in names:
        if not name.endswith(".rels"):
            continue
        try:
            root = ET.fromstring(archive.read(name))
            for element in root.iter():
                if element.tag.rsplit("}", 1)[-1] != "Relationship":
                    continue
                records.append(
                    {
                        "part": name,
                        "id": element.attrib.get("Id", ""),
                        "type": element.attrib.get("Type", ""),
                        "target": element.attrib.get("Target", ""),
                        "target_mode": element.attrib.get("TargetMode", "Internal"),
                    }
                )
        except (ET.ParseError, KeyError) as exc:
            malformed.append({"part": name, "error": type(exc).__name__})
    records.sort(key=lambda item: (item["part"], item["id"], item["target"]))
    return records, malformed


def package_inventory(path: Path, limit: int) -> dict[str, object]:
    with zipfile.ZipFile(path) as archive:
        names = sorted(info.filename for info in archive.infolist() if not info.is_dir())
        rels, malformed = relationships(archive, names)
        content_types = archive.read("[Content_Types].xml")
        risks: dict[str, dict[str, object]] = {}
        for category, pattern in RISK_PATTERNS:
            matched = [name for name in names if pattern.search(name)]
            if matched:
                risks[category] = bounded(matched, limit)
        opaque_names = sorted(
            name
            for name in names
            if re.search(r"vbaProject|^ppt/embeddings/|^_xmlsignatures/", name, re.I)
        )
        opaque_hashes = {name: sha256_bytes(archive.read(name)) for name in opaque_names}
        media = []
        for name in names:
            if name.startswith("ppt/media/"):
                data = archive.read(name)
                media.append({"part": name, "size": len(data), "sha256": sha256_bytes(data)})
        timing_parts = []
        transition_parts = []
        action_parts = []
        active_content: list[dict[str, str]] = []
        hidden_slide_parts = []
        for name in names:
            if not (name.startswith("ppt/slides/") and name.endswith(".xml")):
                continue
            data = archive.read(name)
            try:
                root = ET.fromstring(data)
            except ET.ParseError:
                continue
            local_names = {element.tag.rsplit("}", 1)[-1] for element in root.iter()}
            if "timing" in local_names:
                timing_parts.append(name)
            if "transition" in local_names:
                transition_parts.append(name)
            if local_names.intersection({"hlinkClick", "hlinkHover", "snd", "oleObj"}):
                action_parts.append(name)
            for element in root.iter():
                local_name = element.tag.rsplit("}", 1)[-1]
                if local_name == "oleObj":
                    active_content.append({"kind": "OLE_OBJECT", "part": name, "target": ""})
                if local_name not in {"hlinkClick", "hlinkHover"}:
                    continue
                action = next(
                    (
                        value
                        for key, value in element.attrib.items()
                        if key.rsplit("}", 1)[-1] == "action"
                    ),
                    "",
                )
                action_lower = action.lower()
                kind = "ACTION"
                if "macro" in action_lower:
                    kind = "RUN_MACRO"
                elif "program" in action_lower:
                    kind = "RUN_PROGRAM"
                elif "hlinkfile" in action_lower or "openfile" in action_lower:
                    kind = "OPEN_FILE"
                elif "ole" in action_lower:
                    kind = "OLE_VERB"
                elif "hlinksldjump" in action_lower or "slidejump" in action_lower:
                    kind = "INTERNAL_SLIDE_HYPERLINK"
                active_content.append({"kind": kind, "part": name, "target": action})
            if root.attrib.get("show") in {"0", "false", "False"}:
                hidden_slide_parts.append(name)
        external = [record for record in rels if record["target_mode"].lower() == "external"]
        for relationship in rels:
            relationship_type = relationship["type"].lower()
            source = relationship["part"]
            target = relationship["target"]
            if relationship["target_mode"].lower() == "external":
                kind = (
                    "EXTERNAL_HYPERLINK"
                    if relationship_type.endswith("/hyperlink")
                    else "EXTERNAL_RELATIONSHIP"
                )
                active_content.append({"kind": kind, "part": source, "target": target})
            elif source.startswith("ppt/slides/_rels/") and relationship_type.endswith("/slide"):
                active_content.append(
                    {"kind": "INTERNAL_SLIDE_HYPERLINK", "part": source, "target": target}
                )
        if "ppt/vbaProject.bin" in names:
            active_content.append(
                {"kind": "VBA_PROJECT", "part": "ppt/vbaProject.bin", "target": ""}
            )
        for name in names:
            if name.startswith("ppt/embeddings/"):
                active_content.append({"kind": "EMBEDDED_OBJECT", "part": name, "target": ""})
        active_content.sort(key=lambda item: (item["kind"], item["part"], item["target"]))
        unknown = [name for name in names if not name.startswith(KNOWN_PREFIXES)]
        return {
            "parts": bounded(names, limit),
            "relationships": bounded(rels, limit),
            "external_relationships": bounded(external, limit),
            "active_content": bounded(active_content, limit),
            "risk_parts": risks,
            "unknown_parts": bounded(unknown, limit),
            "malformed_relationships": malformed,
            "media": bounded(media, limit),
            "opaque_hashes": opaque_hashes,
            "macro_enabled_content_type": MACRO_CONTENT_TYPE.encode() in content_types,
            "vba_present": "ppt/vbaProject.bin" in names,
            "vba_sha256": opaque_hashes.get("ppt/vbaProject.bin"),
            "timing_parts": timing_parts,
            "transition_parts": transition_parts,
            "action_parts": action_parts,
            "hidden_slide_parts": hidden_slide_parts,
        }


def run_record(run: object) -> dict[str, object]:
    font = run.font
    color = None
    try:
        color_type = font.color.type
        if color_type is not None:
            if getattr(font.color, "rgb", None) is not None:
                color = str(font.color.rgb)
            elif getattr(font.color, "theme_color", None) is not None:
                color = enum_name(font.color.theme_color)
    except (AttributeError, ValueError):
        color = None
    hyperlink = None
    try:
        hyperlink = run.hyperlink.address
    except (AttributeError, KeyError, ValueError):
        hyperlink = None
    return {
        "text": run.text,
        "font_name": font.name,
        "font_size": font.size,
        "bold": font.bold,
        "italic": font.italic,
        "color": color,
        "hyperlink": hyperlink,
    }


def shape_record(shape: object, index: int) -> dict[str, object]:
    record: dict[str, object] = {
        "index": index,
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": enum_name(shape.shape_type),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
        "rotation": getattr(shape, "rotation", 0.0),
        "is_placeholder": bool(shape.is_placeholder),
    }
    if shape.is_placeholder:
        record["placeholder"] = {
            "idx": shape.placeholder_format.idx,
            "type": enum_name(shape.placeholder_format.type),
        }
    accessibility = {"title": None, "description": None, "decorative": None}
    try:
        for element in shape._element.iter():
            local_name = element.tag.rsplit("}", 1)[-1]
            if local_name == "cNvPr":
                accessibility["title"] = element.attrib.get("title")
                accessibility["description"] = element.attrib.get("descr")
                break
        for element in shape._element.iter():
            if element.tag.rsplit("}", 1)[-1] == "decorative":
                value = element.attrib.get("val", "1")
                accessibility["decorative"] = value not in {"0", "false", "False"}
                break
    except (AttributeError, TypeError):
        pass
    record["accessibility"] = accessibility
    if getattr(shape, "has_text_frame", False):
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            paragraphs.append(
                {
                    "text": paragraph.text,
                    "level": paragraph.level,
                    "alignment": enum_name(paragraph.alignment) if paragraph.alignment else None,
                    "runs": [run_record(run) for run in paragraph.runs],
                }
            )
        record["text"] = shape.text
        record["paragraphs"] = paragraphs
    if getattr(shape, "has_table", False):
        table = shape.table
        record["table"] = {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "cells": [[cell.text for cell in row.cells] for row in table.rows],
        }
    if getattr(shape, "has_chart", False):
        chart = shape.chart
        series = []
        for item in chart.series:
            series.append({"name": item.name, "point_count": len(item.points)})
        record["chart"] = {"chart_type": enum_name(chart.chart_type), "series": series}
    try:
        image = shape.image
        record["image"] = {
            "filename": image.filename,
            "size": image.size,
            "sha256": sha256_bytes(image.blob),
        }
        for attribute in ("crop_left", "crop_top", "crop_right", "crop_bottom"):
            record["image"][attribute] = getattr(shape, attribute, None)
    except (AttributeError, ValueError):
        pass
    return record


def semantic_inventory(path: Path, limit: int) -> dict[str, object]:
    from pptx import Presentation

    presentation = Presentation(path)
    slides = []
    all_fonts: set[str] = set()
    notes_count = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        shapes = [shape_record(shape, index) for index, shape in enumerate(slide.shapes)]
        for shape in shapes:
            for paragraph in shape.get("paragraphs", []):
                for run in paragraph["runs"]:
                    if run["font_name"]:
                        all_fonts.add(str(run["font_name"]))
        title = slide.shapes.title.text if slide.shapes.title is not None else None
        notes_text = None
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text if notes_frame is not None else None
        except (AttributeError, ValueError):
            notes_text = None
        notes_present = bool(notes_text and notes_text.strip())
        notes_count += int(notes_present)
        slides.append(
            {
                "number": slide_index,
                "slide_id": slide.slide_id,
                "layout": slide.slide_layout.name,
                "title": title,
                "shape_count": len(shapes),
                "shapes": bounded(shapes, limit),
                "notes_present": notes_present,
                "notes_text": notes_text,
            }
        )
    return {
        "slide_width": presentation.slide_width,
        "slide_height": presentation.slide_height,
        "slide_count": len(slides),
        "slides": slides,
        "master_count": len(presentation.slide_masters),
        "layout_names": [layout.name for layout in presentation.slide_layouts],
        "fonts": sorted(all_fonts),
        "slides_with_notes": notes_count,
        "core_properties": {
            key: getattr(presentation.core_properties, key, None)
            for key in ("title", "subject", "author", "keywords", "comments", "category")
        },
    }


def inspect(path: Path, *, limit: int = DEFAULT_DETAILS_LIMIT) -> dict[str, object]:
    package = package_inventory(path, limit)
    semantic = semantic_inventory(path, limit)
    risk_count = sum(int(value["count"]) for value in package["risk_parts"].values())
    return {
        "schema_version": SCHEMA_VERSION,
        "ok": True,
        "path": str(path),
        "extension": path.suffix.lower(),
        "size": path.stat().st_size,
        "package": package,
        "presentation": semantic,
        "preservation": {
            "risk_part_count": risk_count,
            "unknown_part_count": package["unknown_parts"]["count"],
            "malformed_relationship_count": len(package["malformed_relationships"]),
            "requires_review": bool(
                risk_count
                or package["unknown_parts"]["count"]
                or package["malformed_relationships"]
            ),
            "note": "Inventory is not proof of python-pptx preservation or runtime behavior.",
        },
    }


def parser() -> argparse.ArgumentParser:
    result = argparse.ArgumentParser(description=__doc__)
    result.add_argument("presentation", type=Path)
    result.add_argument("--details-limit", type=int, default=DEFAULT_DETAILS_LIMIT)
    return result


def main(argv: list[str] | None = None) -> int:
    args = parser().parse_args(argv)
    if args.details_limit < 0:
        print(
            json.dumps(
                {
                    "schema_version": SCHEMA_VERSION,
                    "ok": False,
                    "error": "details-limit must be non-negative",
                }
            )
        )
        return 2
    try:
        payload = inspect(args.presentation.resolve(), limit=args.details_limit)
    except Exception as exc:  # deterministic boundary for library/package parse failures
        payload = {
            "schema_version": SCHEMA_VERSION,
            "ok": False,
            "error": f"{type(exc).__name__}: {exc}",
        }
        print(json.dumps(payload, indent=2, sort_keys=True))
        return 2
    print(json.dumps(payload, indent=2, sort_keys=True, default=str))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
